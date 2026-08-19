"""
Part 1 PPT 生成器：ASABE 2025 经验分享
17 页，简洁通用风格，含流程图/框图。
叙事主线：任务拆解（解耦 → 选型 → 耦合）+ 迭代思路 + 设计细节 + 心得。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os

ASSETS = r"c:\projects\机构设计\assets\asabe"
VIDEO = r"c:\projects\机构设计\assets\video\competition-demo.mp4"
OUT = r"c:\projects\机构设计\part1_asabe\slides.pptx"

# 颜色（与 HTML 端一致）
ACCENT = RGBColor(0xc8, 0x44, 0x2a)   # 主红
BLUE   = RGBColor(0x2a, 0x5c, 0x8a)   # 蓝
GOOD   = RGBColor(0x2a, 0x7c, 0x4a)   # 绿
BAD    = RGBColor(0xa1, 0x30, 0x30)   # 暗红
INK    = RGBColor(0x22, 0x22, 0x22)
INK_SOFT = RGBColor(0x66, 0x66, 0x66)
BG     = RGBColor(0xfa, 0xfa, 0xf7)
PANEL  = RGBColor(0xee, 0xee, 0xea)
LINE   = RGBColor(0xcc, 0xcc, 0xcc)
WHITE  = RGBColor(0xff, 0xff, 0xff)

FONT = "Microsoft YaHei"  # 中文友好

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# === 工具函数 ===

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=WHITE):
    """整页背景填充"""
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect

def title_bar(s, title, sub=None):
    """页顶标题栏"""
    bar = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6))
    bar.text_frame.word_wrap = True
    p = bar.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(32); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = INK
    # 红色短线（装饰）
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.88), Inches(0.7), Emu(45720))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.shadow.inherit = False
    if sub:
        sb = s.shapes.add_textbox(Inches(1.3), Inches(0.92), Inches(11.5), Inches(0.32))
        p2 = sb.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(16); r2.font.name = FONT
        r2.font.bold = True
        r2.font.color.rgb = INK_SOFT
    return bar

def page_number(s, n):
    tb = s.shapes.add_textbox(Inches(12.6), Inches(7.0), Inches(0.7), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{n} / 17"
    r.font.size = Pt(12); r.font.color.rgb = INK_SOFT; r.font.name = FONT

def textbox(s, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = font; r.font.color.rgb = color
    return tb

def bullets(s, x, y, w, h, items, size=17, color=INK):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "• " + item
        r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = color
        p.space_after = Pt(5)
    return tb

def box(s, x, y, w, h, fill=None, line_color=None, line_w=1.0):
    """圆角矩形"""
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill:
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
    else:
        rect.fill.background()
    if line_color:
        rect.line.color.color = line_color
        rect.line.width = Pt(line_w)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    # 调整圆角半径
    rect.adjustments[0] = 0.1
    return rect

def arrow(s, x1, y1, x2, y2, color=INK_SOFT, width=1.5):
    """连接线（带箭头）"""
    conn = s.shapes.add_connector(2, x1, y1, x2, y2)  # 2 = straight
    conn.line.color.color = color
    conn.line.width = Pt(width)
    # 添加箭头
    line_elem = conn.line._get_or_add_ln()
    tailEnd = etree.SubElement(line_elem, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med'); tailEnd.set('h', 'med')
    return conn

def card(s, x, y, w, h, title, body, color=ACCENT, title_size=16, body_size=13):
    """卡片：标题 + 正文"""
    box(s, x, y, w, h, fill=WHITE, line_color=color, line_w=1.5)
    # 左侧色条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(38100), h)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.shadow.inherit = False
    textbox(s, x + Inches(0.18), y + Inches(0.1), w - Inches(0.25), Inches(0.45),
            title, size=title_size, color=color, bold=True)
    textbox(s, x + Inches(0.18), y + Inches(0.55), w - Inches(0.25), h - Inches(0.6),
            body, size=body_size, color=INK)

def image_fit(s, path, x, y, w, h):
    """按比例缩放图片到指定区域，居中"""
    if not os.path.exists(path):
        textbox(s, x, y, w, h, f"[缺图: {os.path.basename(path)}]", size=17, color=INK_SOFT,
                align=PP_ALIGN.CENTER)
        return None
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        new_w = w
        new_h = int(w / img_ratio)
    else:
        new_h = h
        new_w = int(h * img_ratio)
    new_x = x + (w - new_w) // 2
    new_y = y + (h - new_h) // 2
    return s.shapes.add_picture(path, new_x, new_y, new_w, new_h)

# === Slide 1：标题页 ===
s = slide()
bg(s, WHITE)
# 左侧大色块
left = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), SLIDE_H)
left.line.fill.background()
left.fill.solid(); left.fill.fore_color.rgb = ACCENT
left.shadow.inherit = False

textbox(s, Inches(0.5), Inches(0.6), Inches(4), Inches(0.4),
        "ARC 农业机器人俱乐部", size=14, color=WHITE)

textbox(s, Inches(5.0), Inches(1.5), Inches(8), Inches(0.6),
        "机构设计", size=44, color=INK, bold=True)
textbox(s, Inches(5.0), Inches(2.4), Inches(8), Inches(0.5),
        "——从鸡蛋分拣任务讲起", size=22, color=INK_SOFT)

# 分享人（紧接副标题下方）
textbox(s, Inches(5.0), Inches(3.2), Inches(8), Inches(0.4),
        "分享人：宋在扬", size=18, color=INK)
textbox(s, Inches(5.0), Inches(3.7), Inches(8), Inches(0.4),
        "ASABE 2025 国际大学生机器人大赛 · 标准组", size=14, color=INK_SOFT)

# 底部一行话
textbox(s, Inches(5.0), Inches(6.5), Inches(8), Inches(0.5),
        "给 ARC 大一新成员的入门课 · Part 1", size=14, color=INK_SOFT, align=PP_ALIGN.LEFT)

page_number(s, 1)

# === Slide 2：ASABE 2025 任务（先举具体例子） ===
s = slide(); bg(s)
title_bar(s, "ASABE 2025 任务：5 分钟收蛋分拣", "理解任务约束是机构设计的第一步")

# 4 个约束卡片
constraints = [
    ("⏱", "5 min", "限时完成所有动作", ACCENT),
    ("⬚", "12″立方", "机器人整体尺寸约束", BLUE),
    ("🤖", "自主运行", "启动后无需人工干预", GOOD),
    ("🥚", "好坏蛋分类", "好蛋入筐 / 坏蛋排除", BAD),
]
for i, (icon, key, desc, col) in enumerate(constraints):
    x = Inches(0.5 + i * 3.15)
    box(s, x, Inches(1.4), Inches(3.0), Inches(2.4), fill=WHITE, line_color=col, line_w=1.5)
    textbox(s, x, Inches(1.55), Inches(3.0), Inches(0.6),
            icon, size=32, color=col, align=PP_ALIGN.CENTER)
    textbox(s, x, Inches(2.2), Inches(3.0), Inches(0.5),
            key, size=18, color=col, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.2), Inches(2.8), Inches(2.6), Inches(0.9),
            desc, size=14, color=INK, align=PP_ALIGN.CENTER)

# 下方得分
textbox(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
        "得分逻辑：好蛋正确入筐 +坏蛋正确识别 + 自主完成度 + 报告分",
        size=17, color=INK, bold=True)

# 提示框
card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
     "对机构设计的隐含要求",
     "• 高效：5 min 内完成 → 机构要快，不能慢悠悠一个一个抓\n"
     "• 可靠：自主运行 → 机构要稳，不能依赖人工调整\n"
     "• 紧凑：12″立方 → 所有机构要塞进有限空间\n"
     "• 精准：好坏蛋分类 → 末端执行器要带感知（这里是压力传感器）",
     color=BLUE, title_size=16, body_size=12)

page_number(s, 2)

# === Slide 3：赛场与任务（场地 + 得分占比） ===
s = slide(); bg(s)
title_bar(s, "赛场与任务：场地、四项任务与得分", "规则拆解 —— 5 分钟收集分类（35%）是重中之重")

# 左侧：场地示意图
left_x = Inches(0.5); left_w = Inches(6.0)
box(s, left_x, Inches(1.35), left_w, Inches(4.7), fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "arena-layout-standard.png"),
          left_x + Inches(0.15), Inches(1.5), left_w - Inches(0.3), Inches(4.1))
textbox(s, left_x, Inches(5.6), left_w, Inches(0.4),
        "标准组竞技场：8′×8′ · 白线导航 · 红圈为蛋位",
        size=13, color=INK_SOFT, align=PP_ALIGN.CENTER)

# 左下：场地关键尺寸
dims = [
    ("场地", "96″×96″（内 93″×93″）"),
    ("围板高", "3.5″"),
    ("蛋箱", "高 1″，好蛋/坏蛋两个箱"),
    ("限时", "5 分钟 / 主赛一次"),
]
for i, (k, v) in enumerate(dims):
    y = Inches(6.12 + i * 0.3)
    textbox(s, left_x + Inches(0.3), y, Inches(2.0), Inches(0.3),
            k, size=13, color=ACCENT, bold=True)
    textbox(s, left_x + Inches(2.3), y, Inches(3.6), Inches(0.3),
            v, size=13, color=INK)

# 右侧：四项任务 + 得分占比
right_x = Inches(6.8); right_w = Inches(6.0)
textbox(s, right_x, Inches(1.35), right_w, Inches(0.4),
        "四项任务与得分占比", size=17, color=ACCENT, bold=True)

tasks = [
    ("导航", "20%", "沿白线行进，2 分钟限时，按距离/时间计分", BLUE),
    ("蛋识别", "20%", "识别蛋的类型，3 分钟，正确 1 个 1 分", BLUE),
    ("蛋搬运", "25%", "磁性蛋移入收集箱，3 分钟，蛋易碎", ACCENT),
    ("收集+分类", "35%", "5 分钟主赛：收集并分类，好/坏蛋分箱", BAD),
]
for i, (name, pct, desc, col) in enumerate(tasks):
    y = Inches(1.9 + i * 1.05)
    box(s, right_x, y, right_w, Inches(0.9), fill=WHITE, line_color=col, line_w=1.2)
    # 名称
    textbox(s, right_x + Inches(0.15), y + Inches(0.08), Inches(1.7), Inches(0.4),
            name, size=15, color=col, bold=True)
    # 占比
    textbox(s, right_x + Inches(1.75), y + Inches(0.08), Inches(0.9), Inches(0.4),
            pct, size=16, color=col, bold=True)
    # 要求
    textbox(s, right_x + Inches(2.7), y + Inches(0.15), Inches(3.2), Inches(0.7),
            desc, size=12, color=INK)

# 底部：总分公式提示
card(s, Inches(6.8), Inches(6.15), Inches(6.0), Inches(1.0),
     "总分 = 性能分 × 报告分",
     "收集分类占性能分 35% —— 我们的机构设计就是围绕这一题展开的",
     color=BAD, title_size=14, body_size=12)

page_number(s, 3)

# === Slide 4：任务拆解方法论（再提炼方法论） ===
s = slide(); bg(s)
title_bar(s, "方法论 · 任务拆解：解耦 → 选型 → 耦合", "Part 1 的核心思想，贯穿整节课")

# 上半部：解耦
textbox(s, Inches(0.5), Inches(1.3), Inches(12.5), Inches(0.4),
        "① 解耦（Decouple）：把复杂任务拆成可独立思考的子问题",
        size=16, color=ACCENT, bold=True)

# 流程图：任务居中在上，子问题在下方左右分布——箭头不穿过任何框
# 任务框（居中：x = (13.333 - 3.2)/2 ≈ 5.07）
task_x, task_y, task_w, task_h = Inches(5.07), Inches(1.95), Inches(3.2), Inches(0.8)
box(s, task_x, task_y, task_w, task_h, fill=ACCENT)
textbox(s, task_x, task_y + Inches(0.05), task_w, task_h,
        "鸡蛋分拣任务", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 子问题 A（左下）和 B（右下）
sub_a_x = Inches(1.5); sub_b_x = Inches(8.0)
sub_y, sub_w, sub_h = Inches(3.3), Inches(3.8), Inches(0.8)
box(s, sub_a_x, sub_y, sub_w, sub_h, line_color=BLUE, line_w=1.5)
textbox(s, sub_a_x, sub_y + Inches(0.05), sub_w, sub_h,
        "子问题 A：怎么收集？", size=17, color=INK, align=PP_ALIGN.CENTER)
box(s, sub_b_x, sub_y, sub_w, sub_h, line_color=BLUE, line_w=1.5)
textbox(s, sub_b_x, sub_y + Inches(0.05), sub_w, sub_h,
        "子问题 B：怎么分拣？", size=17, color=INK, align=PP_ALIGN.CENTER)

# 箭头：T 形分叉（从上往下）
#   [任务]
#     |
# +---+---+
# |       |
#[A]     [B]
# 主干（任务底部中心 → 横梁中心，无箭头）
seg = s.shapes.add_connector(1, Inches(6.67), Inches(2.75), Inches(6.67), Inches(3.05))
seg.line.color.color = INK_SOFT; seg.line.width = Pt(1.5)
# 横梁（无箭头）
seg2 = s.shapes.add_connector(1, Inches(3.4), Inches(3.05), Inches(9.9), Inches(3.05))
seg2.line.color.color = INK_SOFT; seg2.line.width = Pt(1.5)
# 左下：横梁 → A 顶部中心（带箭头，垂直向下）
arrow(s, Inches(3.4), Inches(3.05), Inches(3.4), Inches(3.3))
# 右下：横梁 → B 顶部中心（带箭头，垂直向下）
arrow(s, Inches(9.9), Inches(3.05), Inches(9.9), Inches(3.3))

# 下半部：选型 + 耦合（往下移，避免与子问题框重叠）
textbox(s, Inches(0.5), Inches(4.45), Inches(12.5), Inches(0.4),
        "② 选型 + 耦合：各子问题分别选型，再考虑联动合并",
        size=16, color=ACCENT, bold=True)

coupling = box(s, Inches(3.5), Inches(5.05), Inches(6.0), Inches(0.8), fill=PANEL, line_color=ACCENT, line_w=1.5)
textbox(s, Inches(3.5), Inches(5.1), Inches(6.0), Inches(0.8),
        "关键问题：A 和 B 能不能合并到一个机构里实现？",
        size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)

# 结论
card(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.2),
     "我们的答案（后面会展开）",
     "收集用「刮板扫入」，分拣用「压力传感器称重」——\n两者合并到「刮板 + 压力传感器一体化」结构里，一个机构同时完成两个子任务。",
     color=ACCENT, title_size=18, body_size=12)

page_number(s, 4)

# === Slide 5：收集方案对比表 ===
s = slide(); bg(s)
title_bar(s, "解耦-收集方案：5 选 1", "穷举候选方案 → 用约束筛选")

# 表头
cols = [("方案", 2.2, ACCENT), ("原理", 3.2, None), ("优点", 2.7, None),
        ("缺点", 2.7, None), ("是否选用", 1.5, None)]
rows = [
    ("机械臂抓", "多自由度臂 + 夹爪", "通用、灵活", "慢、控制复杂", ""),
    ("真空吸盘", "负压吸附", "对鸡蛋友好", "需气源、易漏", ""),
    ("网兜捡球式", "皮筋网兜 + 下压", "结构简单、对蛋友好", "需下压动作、单向", ""),
    ("刮板扫入", "侧向单方向推", "结构最简、可一体化称重", "只能特定方向", "★ 我们选用"),
    ("履带拾取", "传送带底部收集", "可大批量", "体积大、易卡", ""),
]
table_top = Inches(1.4)
row_h = Inches(0.7)
header_h = Inches(0.5)

# 表头
x = Inches(0.5)
for name, w, col in cols:
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, Inches(w), header_h)
    rect.fill.solid(); rect.fill.fore_color.rgb = ACCENT
    rect.line.color.color = WHITE; rect.line.width = Pt(1)
    rect.shadow.inherit = False
    textbox(s, x, table_top, Inches(w), header_h,
            name, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += Inches(w)

# 数据行
for ri, row in enumerate(rows):
    y = table_top + header_h + Inches(0.7 * ri)
    x = Inches(0.5)
    chosen = row[4].startswith("★")
    for ci, (cell, (_, w, _)) in enumerate(zip(row, cols)):
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(w), row_h)
        if chosen:
            rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0xfd, 0xe7, 0xe0)
        else:
            rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
        rect.line.color.color = LINE; rect.line.width = Pt(0.5)
        rect.shadow.inherit = False
        col = ACCENT if chosen else INK
        bold = chosen and ci == 4
        textbox(s, x + Inches(0.1), y + Inches(0.05), Inches(w - 0.2), Inches(0.6),
                cell, size=17, color=col, bold=bold, align=PP_ALIGN.CENTER)
        x += Inches(w)

# 底部决策说明
y_dec = table_top + header_h + Inches(0.7 * 5) + Inches(0.25)
card(s, Inches(0.5), y_dec, Inches(12.3), Inches(1.4),
     "为什么选刮板？",
     "① 结构最简（一个舵机 + 一块板）→ 满足 12″立方约束\n"
     "② 可与压力传感器一体化 → 满足「收集 + 分拣」耦合需求\n"
     "③ 鸡蛋易碎，不夹比夹更安全",
     color=ACCENT, title_size=16, body_size=12)

page_number(s, 5)

# === Slide 6：分拣方案对比表 ===
s = slide(); bg(s)
title_bar(s, "解耦-分拣方案：5 选 1", "分拣 = 区分好蛋/坏蛋")

cols2 = [("方案", 2.2, ACCENT), ("原理", 3.2, None), ("优点", 2.7, None),
         ("缺点", 2.7, None), ("是否选用", 1.5, None)]
rows2 = [
    ("视觉识别", "摄像头 + 图像分类", "信息量大、可识别裂纹", "算力高、训练复杂", ""),
    ("颜色识别", "颜色传感器", "便宜、可靠", "对灯光敏感", "他组选用"),
    ("重量检测", "压力传感器", "直接、可靠", "需接触", "★ 我们选用"),
    ("丝线分拣", "不同紧度丝线", "无传感器、自然分拣", "需精细调参数", "他组选用"),
    ("大小检测", "光电测距", "结构简单", "好坏蛋大小差不多", ""),
]
table_top = Inches(1.4)
row_h = Inches(0.7)
header_h = Inches(0.5)

x = Inches(0.5)
for name, w, col in cols2:
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, Inches(w), header_h)
    rect.fill.solid(); rect.fill.fore_color.rgb = ACCENT
    rect.line.color.color = WHITE; rect.line.width = Pt(1)
    rect.shadow.inherit = False
    textbox(s, x, table_top, Inches(w), header_h,
            name, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += Inches(w)

for ri, row in enumerate(rows2):
    y = table_top + header_h + Inches(0.7 * ri)
    x = Inches(0.5)
    chosen = row[4].startswith("★")
    other = row[4].startswith("他组")
    for ci, (cell, (_, w, _)) in enumerate(zip(row, cols2)):
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(w), row_h)
        if chosen:
            rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0xfd, 0xe7, 0xe0)
        else:
            rect.fill.solid(); rect.fill.fore_color.rgb = WHITE
        rect.line.color.color = LINE; rect.line.width = Pt(0.5)
        rect.shadow.inherit = False
        col = ACCENT if chosen else INK
        bold = (chosen or other) and ci == 4
        textbox(s, x + Inches(0.1), y + Inches(0.05), Inches(w - 0.2), Inches(0.6),
                cell, size=17, color=col, bold=bold, align=PP_ALIGN.CENTER)
        x += Inches(w)

y_dec = table_top + header_h + Inches(0.7 * 5) + Inches(0.25)
card(s, Inches(0.5), y_dec, Inches(12.3), Inches(1.4),
     "为什么选重量检测？",
     "① 好蛋/坏蛋（破损、变质）重量差异明显 → 物理特征可靠\n"
     "② 可与刮板结构耦合 → 刮板底部贴压力传感器，扫入即称重\n"
     "③ 不需要额外传感器/算力 → 满足「自主 + 紧凑」约束\n"
     "注：丝线分拣是同赛场另一组的方案——重的蛋穿过丝线落下，轻的弹开，构思很巧。",
     color=ACCENT, title_size=16, body_size=11)

page_number(s, 6)

# === Slide 7：耦合方案对比（不止一种解法） ===
s = slide(); bg(s)
title_bar(s, "耦合不止一种：两组方案对照", "同样把「收集 + 分拣」合并，思路可以很不一样")

# 左栏：友组方案（蓝）
left_x = Inches(0.5); col_w = Inches(6.0)
box(s, left_x, Inches(1.3), col_w, Inches(0.6), fill=BLUE)
textbox(s, left_x, Inches(1.35), col_w, Inches(0.6),
        "友组 · 网兜拾取 + 丝线分拣",
        size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 两张图左右并排
img_y = Inches(2.05); img_h = Inches(3.0)
half_w = Inches(2.9)
# 左图：采集框
box(s, left_x, img_y, half_w, img_h, fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "friend-collection.png"),
          left_x + Inches(0.1), img_y + Inches(0.1), half_w - Inches(0.2), img_h - Inches(0.4))
textbox(s, left_x, img_y + img_h - Inches(0.3), half_w, Inches(0.3),
        "网兜收集框", size=14, color=INK_SOFT, align=PP_ALIGN.CENTER)
# 右图：分选上视图
img_x2 = left_x + half_w + Inches(0.2)
box(s, img_x2, img_y, half_w, img_h, fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "friend-sorting.png"),
          img_x2 + Inches(0.1), img_y + Inches(0.1), half_w - Inches(0.2), img_h - Inches(0.4))
textbox(s, img_x2, img_y + img_h - Inches(0.3), half_w, Inches(0.3),
        "丝线分选（上视）", size=14, color=INK_SOFT, align=PP_ALIGN.CENTER)

# 友组要点
bullets(s, left_x, Inches(5.2), col_w, Inches(1.5),
        ["收集与分拣用两个独立机构",
         "丝线按紧度分选：重蛋穿线落下，轻蛋弹开",
         "构思巧妙，机构之间互不干扰"],
        size=14)

# 右栏：我们方案（红）
right_x = Inches(6.8)
box(s, right_x, Inches(1.3), col_w, Inches(0.6), fill=ACCENT)
textbox(s, right_x, Inches(1.35), col_w, Inches(0.6),
        "我们 · 刮板 + 压力传感器 一体化",
        size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 用一张已有图代表我们方案
box(s, right_x, Inches(2.05), col_w, Inches(3.0), fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "end-effector-picking.png"),
          right_x + Inches(0.15), Inches(2.2), col_w - Inches(0.3), Inches(2.7))
textbox(s, right_x, Inches(4.75), col_w, Inches(0.3),
        "末端执行器（刮板 + 蜂窝筐 + 压力传感器）",
        size=14, color=INK_SOFT, align=PP_ALIGN.CENTER)

# 我们要点（精简：细节交给第 7 页）
bullets(s, right_x, Inches(5.2), col_w, Inches(1.5),
        ["一个结构件完成「收集 + 分拣」",
         "具体怎么实现 → 下页详述"],
        size=14)

# 底部结论
card(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.8),
     "殊途同归",
     "耦合没有标准答案——关键是找到子问题的「共享点」。下页详述我们的耦合方案。",
     color=ACCENT, title_size=15, body_size=13)

page_number(s, 7)

# === Slide 8：耦合页（我们的方案详述） ===
s = slide(); bg(s)
title_bar(s, "耦合：刮板 + 压力传感器 一体化", "两个子问题的解在这里合并")

# 左侧图（末端执行器实物图）
img_box_x, img_box_y = Inches(0.5), Inches(1.4)
img_box_w, img_box_h = Inches(6.5), Inches(5.5)
box(s, img_box_x, img_box_y, img_box_w, img_box_h, fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "end-effector-picking.png"),
          img_box_x + Inches(0.2), img_box_y + Inches(0.3),
          img_box_w - Inches(0.4), img_box_h - Inches(0.6))
textbox(s, img_box_x, img_box_y + img_box_h - Inches(0.3), img_box_w, Inches(0.3),
        "末端执行器（刮板 + 蜂窝筐 + 压力传感器）",
        size=17, color=INK_SOFT, align=PP_ALIGN.CENTER)

# 右侧说明
right_x = Inches(7.3)
right_w = Inches(5.5)

textbox(s, right_x, Inches(1.4), right_w, Inches(0.4),
        "结构 = 解 A + 解 B 的物理交集",
        size=17, color=ACCENT, bold=True)

# 三个要点
points = [
    ("① 收集", "刮板由舵机驱动，绕轴摆动 → 鸡蛋被扫入蜂窝筐"),
    ("② 分拣", "蜂窝筐底部嵌入压力传感器 → 鸡蛋进入即称重"),
    ("③ 一体化", "同一结构件承担两个功能 → 不需要两个独立机构"),
]
for i, (k, v) in enumerate(points):
    y = Inches(2.0 + i * 1.0)
    box(s, right_x, y, right_w, Inches(0.9), fill=WHITE, line_color=ACCENT, line_w=1.0)
    textbox(s, right_x + Inches(0.15), y + Inches(0.08), Inches(2), Inches(0.4),
            k, size=14, color=ACCENT, bold=True)
    textbox(s, right_x + Inches(0.15), y + Inches(0.42), right_w - Inches(0.3), Inches(0.5),
            v, size=17, color=INK)

# 底部方法论提示
textbox(s, right_x, Inches(5.3), right_w, Inches(0.5),
        "✓ 耦合思维 ≠ 拼接",
        size=14, color=GOOD, bold=True)
textbox(s, right_x, Inches(5.7), right_w, Inches(1.2),
        "不是为了省零件而合并，\n而是因为两个子任务在物理上共享同一个执行点（鸡蛋被夹持的位置）。\n机构设计要找的就是这种「共享点」。",
        size=14, color=INK)

page_number(s, 8)

# === Slide 9：设计细节（末端执行器关键参数） ===
s = slide(); bg(s)
title_bar(s, "设计细节：末端执行器的关键参数", "方案落地后，参数决定成败")

# 左侧图
img_x, img_y = Inches(0.5), Inches(1.35)
img_w, img_h = Inches(5.6), Inches(5.55)
box(s, img_x, img_y, img_w, img_h, fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "end-effector-picking.png"),
          img_x + Inches(0.15), img_y + Inches(0.15), img_w - Inches(0.3), img_h - Inches(0.5))
textbox(s, img_x, img_y + img_h - Inches(0.3), img_w, Inches(0.3),
        "末端执行器：刮板 + 蜂窝筐 + 压力传感器",
        size=14, color=INK_SOFT, align=PP_ALIGN.CENTER)

# 右侧参数卡片（2 列 × 3 行）
param_cards = [
    ("刮板", "薄塑料板 + 微舵机\n边缘限位脊 → 蛋与传感器完全接触", ACCENT),
    ("压力传感器", "薄膜式 40×40mm\n量程 20g–10kg，按阈值分类", ACCENT),
    ("收集筐", "3D 打印 PLA，梯形上窄下宽\n蜂窝镂空 + 加强筋 → 减重保刚度", BLUE),
    ("好蛋释放", "5mm 泡棉内衬斜面\n双舵机铰链翻板", GOOD),
    ("坏蛋释放", "布料张紧变斜坡\n存储 / 释放一体化", BAD),
    ("全机成本", "$273.39\n处理器（Arduino+OpenMV）仅 $53.6", INK),
]
card_x0 = Inches(6.4); card_x1 = Inches(9.85); card_w = Inches(3.15)
for i, (t, body, col) in enumerate(param_cards):
    r, c = divmod(i, 2)
    x = card_x0 if c == 0 else card_x1
    y = Inches(1.35 + r * 1.85)
    box(s, x, y, card_w, Inches(1.7), fill=WHITE, line_color=col, line_w=1.2)
    textbox(s, x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.4),
            t, size=15, color=col, bold=True)
    textbox(s, x + Inches(0.15), y + Inches(0.55), card_w - Inches(0.3), Inches(1.1),
            body, size=13, color=INK)

page_number(s, 9)

# === Slide 10：迭代思路（方案只是起点） ===
s = slide(); bg(s)
title_bar(s, "迭代思路：方案只是起点", "方案确定后，靠一次次迭代把它做完善")

# 点题
textbox(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
        "没有一次成型的设计——车是一版一版迭代出来的：",
        size=16, color=INK, bold=True)

# 5 个迭代条目：问题（左）→ 迭代方案（右）
iterations = [
    ("① 收集筐太重", "实心壁 → 蜂窝镂空 + 加强筋", "减重明显，刚度不变"),
    ("② 磁性蛋易碎", "加柔性泡沫垫缓冲", "破损率下降"),
    ("③ 存储按好:坏=3:1 设计，实际 1:1", "坏蛋仓改纯布料，收集+释放一体化", "存储空间大幅增加"),
    ("④ 激光对不齐鸡蛋", "微调传感器高度，激光正对鸡蛋", "探测稳定"),
    ("⑤ 单层结构电路机构杂糅", "改两层：下层电路板 / 上层机构", "装配、调试、排查都方便"),
]
row_y = Inches(1.85); row_h = Inches(0.92); row_gap = Inches(0.08)
prob_x = Inches(0.5); prob_w = Inches(4.6)
iter_x = Inches(5.55); iter_w = Inches(7.25)
for i, (prob, sol, eff) in enumerate(iterations):
    y = row_y + int(row_h + row_gap) * i
    # 问题框（浅红底）
    box(s, prob_x, y, prob_w, row_h, fill=RGBColor(0xfb, 0xe9, 0xe7), line_color=BAD, line_w=1.0)
    textbox(s, prob_x + Inches(0.15), y + Inches(0.14), prob_w - Inches(0.3), Inches(0.7),
            prob, size=14, color=INK)
    # 箭头
    textbox(s, Inches(5.16), y + Inches(0.22), Inches(0.35), Inches(0.5),
            "→", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    # 迭代框（浅绿底）
    box(s, iter_x, y, iter_w, row_h, fill=RGBColor(0xe9, 0xf3, 0xec), line_color=GOOD, line_w=1.0)
    textbox(s, iter_x + Inches(0.15), y + Inches(0.08), iter_w - Inches(0.3), Inches(0.42),
            sol, size=14, color=GOOD, bold=True)
    textbox(s, iter_x + Inches(0.15), y + Inches(0.5), iter_w - Inches(0.3), Inches(0.4),
            "→ " + eff, size=13, color=INK)

page_number(s, 10)

# === Slide 11：整机 3D 总览 ===
s = slide(); bg(s)
title_bar(s, "整机总览：模块化标注", "5 个子系统，每个都有对应的机构设计点")

# 大图
img_x, img_y = Inches(0.5), Inches(1.3)
img_w, img_h = Inches(8.5), Inches(5.5)
box(s, img_x, img_y, img_w, img_h, fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "robot-full-design.png"),
          img_x + Inches(0.2), img_y + Inches(0.2),
          img_w - Inches(0.4), img_h - Inches(0.4))

# 右侧模块列表
right_x = Inches(9.3); right_w = Inches(3.5)
textbox(s, right_x, Inches(1.3), right_w, Inches(0.4),
        "5 个子系统", size=17, color=ACCENT, bold=True)
modules = [
    ("① 移动", "4× 麦轮 + 42 步进", BLUE),
    ("② 机械臂", "3-DOF 串联关节", ACCENT),
    ("③ 末端执行", "刮板 + 蜂窝筐", ACCENT),
    ("④ 好蛋释放", "倾斜板 + 5mm 泡棉", GOOD),
    ("⑤ 坏蛋释放", "软织物张紧", BAD),
]
for i, (k, v, col) in enumerate(modules):
    y = Inches(1.8 + i * 0.95)
    box(s, right_x, y, right_w, Inches(0.85), fill=WHITE, line_color=col, line_w=1.0)
    textbox(s, right_x + Inches(0.1), y + Inches(0.05), right_w - Inches(0.2), Inches(0.4),
            k, size=14, color=col, bold=True)
    textbox(s, right_x + Inches(0.1), y + Inches(0.42), right_w - Inches(0.2), Inches(0.4),
            v, size=17, color=INK)

page_number(s, 11)

# === Slide 12：释放机构（合并页） ===
s = slide(); bg(s)
title_bar(s, "释放机构：好蛋 + 坏蛋", "两种独立机构，对应两种释放逻辑")

# 左：好蛋释放
left_x = Inches(0.5); left_w = Inches(6.0)
box(s, left_x, Inches(1.3), left_w, Inches(0.6), fill=GOOD)
textbox(s, left_x, Inches(1.35), left_w, Inches(0.6),
        "好蛋释放 · 倾斜板 + 5mm 泡棉",
        size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
box(s, left_x, Inches(2.0), left_w, Inches(3.5), fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "storage-release.png"),
          left_x + Inches(0.15), Inches(2.15), left_w - Inches(0.3), Inches(3.2))
bullets(s, left_x, Inches(5.6), left_w, Inches(1.5),
        ["舵机驱动铰链板（简化四杆）",
         "5mm 泡棉缓冲 → 鸡蛋不破",
         "侧置舵机收紧 → 板倾斜 → 蛋滑出"],
        size=11)

# 右：坏蛋释放
right_x = Inches(6.8); right_w = Inches(6.0)
box(s, right_x, Inches(1.3), right_w, Inches(0.6), fill=BAD)
textbox(s, right_x, Inches(1.35), right_w, Inches(0.6),
        "坏蛋释放 · 软织物张紧",
        size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
box(s, right_x, Inches(2.0), right_w, Inches(3.5), fill=PANEL, line_color=LINE, line_w=0.5)
image_fit(s, os.path.join(ASSETS, "material-fabric.png"),
          right_x + Inches(0.15), Inches(2.15), right_w - Inches(0.3), Inches(3.2))
bullets(s, right_x, Inches(5.6), right_w, Inches(1.5),
        ["软织物平时张紧成斜坡",
         "舵机收紧 → 坡度增加 → 坏蛋滑出",
         "柔顺机构思路（无刚性结构）"],
        size=14)

page_number(s, 12)

# === Slide 13：实战运行闭环 ===
s = slide(); bg(s)
title_bar(s, "实战流程：收蛋 → 称重 → 分类 → 释放", "一次完整的工作循环")

# 纵向 6 步
steps = [
    ("① 寻线移动", "相机识别引导线，麦轮开环行进"),
    ("② 发现鸡蛋", "激光测距探测 20–200mm 内的蛋"),
    ("③ 收蛋", "机械臂降下，刮板扫入蜂窝筐"),
    ("④ 称重分类", "薄膜传感器阈值 → 绿灯好蛋 / 红灯坏蛋"),
    ("⑤ 分区存储", "好蛋仓 + 坏蛋仓，互不干扰"),
    ("⑥ 统一释放", "整场好/坏蛋各释放一次，自动回起点"),
]
flow_x = Inches(1.8); flow_w = Inches(9.7); flow_h = Inches(0.6)
for i, (t, d) in enumerate(steps):
    y = Inches(1.4) + int(Inches(0.82)) * i
    box(s, flow_x, y, flow_w, flow_h, fill=PANEL if i % 2 == 0 else WHITE, line_color=BLUE, line_w=1.0)
    textbox(s, flow_x + Inches(0.25), y + Inches(0.05), Inches(2.5), flow_h,
            t, size=15, color=ACCENT, bold=True)
    textbox(s, flow_x + Inches(2.8), y + Inches(0.08), flow_w - Inches(3.0), flow_h,
            d, size=14, color=INK)
    if i < 5:
        arrow(s, Inches(6.65), y + flow_h, Inches(6.65), y + Inches(0.82), color=INK_SOFT, width=1.5)

# 底部亮点
textbox(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.4),
        "设计亮点：好/坏蛋整场各释放一次 —— 先存储、最后统一倾倒，省时间",
        size=15, color=GOOD, bold=True, align=PP_ALIGN.CENTER)

page_number(s, 13)

# === Slide 14：比赛演示视频 ===
s = slide(); bg(s)
title_bar(s, "比赛演示视频", "整机集成后的实际运行")

# 视频占位
video_x, video_y = Inches(2.0), Inches(1.4)
video_w, video_h = Inches(9.3), Inches(5.2)
box(s, video_x, video_y, video_w, video_h, fill=RGBColor(0x22,0x22,0x22), line_color=ACCENT, line_w=2)

if os.path.exists(VIDEO):
    s.shapes.add_movie(VIDEO, video_x + Inches(0.1), video_y + Inches(0.1),
                       video_w - Inches(0.2), video_h - Inches(0.2),
                       poster_frame_image=None)
else:
    textbox(s, video_x, video_y + video_h/2 - Inches(0.3), video_w, Inches(0.4),
            "[视频文件未找到]", size=18, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, video_x, video_y + video_h/2 + Inches(0.1), video_w, Inches(0.4),
            "assets/video/competition-demo.mp4", size=14, color=WHITE, align=PP_ALIGN.CENTER)

textbox(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
        "关注：自主运行、5 min 内的节奏、收集-称重-释放的完整循环",
        size=15, color=INK_SOFT, align=PP_ALIGN.CENTER)

page_number(s, 14)

# === Slide 15：机构设计心得 ===
s = slide(); bg(s)
title_bar(s, "机构设计心得：实战教会我们的 6 件事", "把经验变成下一次设计的起点")

lessons = [
    ("简洁有效，不是越复杂越好",
     "开环导航省掉编码器/IMU —— 能用一个舵机解决，就不用两个；简单 = 可靠 = 便宜"),
    ("针对最脆弱的环节设计",
     "磁性蛋掉地就散架 → 泡棉垫缓冲；先找最脆弱的点，再设计"),
    ("模块化：并行推进 + 方便迭代",
     "5 人分工互不干扰；迭代只改对应模块，排查故障也快"),
    ("成本是约束也是动力",
     "全机 $273.39，处理器仅 $53.6 —— 约束逼出简单方案"),
    ("用评分规则倒推设计",
     "总分 = 性能 × 报告；测试数据占 5 分，我们的报告恰好缺了（反面教材）"),
    ("考虑环境因素",
     "场地材质、光照影响视觉 → 影响机构执行和整车运动；规则方都改了场地材质（胶合板 → 泡棉垫）"),
]
card_w = Inches(6.05); card_h = Inches(1.75)
xs = [Inches(0.5), Inches(6.8)]
ys = [Inches(1.35), Inches(3.25), Inches(5.15)]
for i, (t, d) in enumerate(lessons):
    r, c = divmod(i, 2)
    x = xs[c]; y = ys[r]
    box(s, x, y, card_w, card_h, fill=WHITE, line_color=ACCENT, line_w=1.2)
    textbox(s, x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.4),
            str(i + 1), size=18, color=ACCENT, bold=True)
    textbox(s, x + Inches(0.65), y + Inches(0.1), card_w - Inches(0.8), Inches(0.4),
            t, size=15, color=ACCENT, bold=True)
    textbox(s, x + Inches(0.65), y + Inches(0.55), card_w - Inches(0.8), Inches(1.1),
            d, size=13, color=INK)

page_number(s, 15)

# === Slide 16：方法论总结（通用流程图） ===
s = slide(); bg(s)
title_bar(s, "方法论回顾：解耦 → 选型 → 耦合", "一套可复用的机构设计思路，不止用于鸡蛋分拣")

# 纵向流程图：5 个阶段框，居中排列
flow_x = Inches(4.5); flow_w = Inches(4.3); flow_h = Inches(0.7)
stages = [
    ("复杂任务",     "一个看起来很难下手的大问题",   PANEL, INK),
    ("① 解耦",       "拆成可独立思考的子问题",       ACCENT, WHITE),
    ("② 选型",       "每个子问题穷举候选 → 用约束筛选", ACCENT, WHITE),
    ("③ 耦合",       "找子问题之间的「共享点」合并",   ACCENT, WHITE),
    ("一体化机构",   "一个结构同时完成多个子任务 ✓",   GOOD,  WHITE),
]

# 框的纵向起始与间距
top_y = Inches(1.35)
gap = Inches(0.95)  # 框高 0.7 + 箭头空间 0.25

for i, (name, desc, fill_col, txt_col) in enumerate(stages):
    y = top_y + int(gap) * i
    # 阶段框
    is_start = (i == 0)
    is_end = (i == len(stages) - 1)
    if is_start or is_end:
        box(s, flow_x, y, flow_w, flow_h, fill=fill_col, line_color=fill_col, line_w=1.5)
    else:
        box(s, flow_x, y, flow_w, flow_h, fill=fill_col, line_color=fill_col, line_w=1.5)
    textbox(s, flow_x, y + Inches(0.05), flow_w, flow_h,
            name, size=18, color=txt_col, bold=True, align=PP_ALIGN.CENTER)
    # 右侧说明
    textbox(s, flow_x + flow_w + Inches(0.3), y + Inches(0.12), Inches(4.0), Inches(0.5),
            desc, size=15, color=INK_SOFT)
    # 除最后一个外，画向下箭头
    if not is_end:
        arrow_x = flow_x + flow_w // 2
        arrow(s, arrow_x, y + flow_h, arrow_x, y + gap, color=INK_SOFT, width=1.75)

# 左侧留白区：用一句话点题
textbox(s, Inches(0.5), Inches(2.4), Inches(3.7), Inches(2.0),
        "解耦让问题\n变简单\n\n选型让方案\n有依据\n\n耦合让结构\n更精炼",
        size=15, color=INK_SOFT, align=PP_ALIGN.CENTER)

page_number(s, 16)

# === Slide 17：桥段 + Q&A ===
s = slide(); bg(s)
title_bar(s, "桥段回顾 → 进入 Part 2", "下面是 HTML 互动页面，可以现场调参数")

# 桥段列表
textbox(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
        "Part 1 提到的 → Part 2 对应页面：",
        size=17, color=ACCENT, bold=True)

bridges = [
    ("3-DOF 机械臂", "自由度页 + 复杂曲线页", "dof.html"),
    ("刮板扫入", "直线运动页 + 夹持拾取页", "motion-gripping.html"),
    ("铰链板释放", "摆动页（四杆机构 hero）", "motion-oscillation.html"),
    ("步进电机直驱", "间歇运动页 + 减速传动页", "motion-intermittent.html"),
]
for i, (p1, p2, link) in enumerate(bridges):
    y = Inches(1.85 + i * 0.6)
    # Part 1
    box(s, Inches(0.5), y, Inches(4.5), Inches(0.5), fill=PANEL, line_color=ACCENT, line_w=1.0)
    textbox(s, Inches(0.5), y + Inches(0.1), Inches(4.5), Inches(0.4),
            p1, size=14, color=INK, align=PP_ALIGN.CENTER)
    # 箭头
    textbox(s, Inches(5.0), y + Inches(0.05), Inches(0.5), Inches(0.4),
            "→", size=18, color=INK_SOFT, align=PP_ALIGN.CENTER)
    # Part 2
    box(s, Inches(5.5), y, Inches(7.3), Inches(0.5), fill=WHITE, line_color=BLUE, line_w=1.0)
    textbox(s, Inches(5.5), y + Inches(0.1), Inches(7.3), Inches(0.4),
            p2, size=14, color=BLUE, align=PP_ALIGN.CENTER)

# Q&A 区
qa_y = Inches(5.0)
box(s, Inches(0.5), qa_y, Inches(12.3), Inches(2.0), fill=PANEL, line_color=ACCENT, line_w=2.0)
textbox(s, Inches(0.8), qa_y + Inches(0.15), Inches(11.5), Inches(0.5),
        "Q & A", size=24, color=ACCENT, bold=True)
textbox(s, Inches(0.8), qa_y + Inches(0.7), Inches(11.5), Inches(1.3),
        "现在欢迎提问。问题方向：\n"
        "• 任何机构的选择理由\n"
        "• 实物制作中的具体问题（材料、加工、装配）\n"
        "• ASABE 比赛相关",
        size=15, color=INK)

page_number(s, 17)

# === 保存 ===
prs.save(OUT)
print(f"OK: {OUT}")
print(f"共 {len(prs.slides)} 页")
